"""Template PPTX parser for extracting theme and layouts"""
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE, PP_PLACEHOLDER
from typing import Dict, List, Any
import uuid
from pathlib import Path

class TemplateParser:
    """Parse PPTX template to extract theme metadata and layout catalog"""
    
    def __init__(self, file_path: Path):
        self.prs = Presentation(str(file_path))
        self.template_id = str(uuid.uuid4())
    
    def parse(self) -> Dict[str, Any]:
        """Extract theme and layouts from template"""
        theme_meta = self._extract_theme()
        layout_catalog = self._extract_layouts()
        
        return {
            "template_id": self.template_id,
            "theme_meta": theme_meta,
            "layout_catalog": layout_catalog
        }
    
    def _extract_theme(self) -> Dict[str, Any]:
        """Extract theme colors and fonts"""
        theme = {
            "fonts": {
                "title": "Calibri",
                "body": "Calibri"
            },
            "colors": {},
            "page_size": {
                "w": self.prs.slide_width,
                "h": self.prs.slide_height
            }
        }
        
        # Try to get theme fonts
        try:
            if self.prs.slide_master.theme:
                theme_element = self.prs.slide_master.theme.element
                
                # Extract major (title) and minor (body) fonts
                font_scheme = theme_element.find('.//a:fontScheme', 
                    {'a': 'http://schemas.openxmlformats.org/drawingml/2006/main'})
                if font_scheme is not None:
                    major_font = font_scheme.find('.//a:majorFont/a:latin',
                        {'a': 'http://schemas.openxmlformats.org/drawingml/2006/main'})
                    minor_font = font_scheme.find('.//a:minorFont/a:latin',
                        {'a': 'http://schemas.openxmlformats.org/drawingml/2006/main'})
                    
                    if major_font is not None and 'typeface' in major_font.attrib:
                        theme["fonts"]["title"] = major_font.attrib['typeface']
                    if minor_font is not None and 'typeface' in minor_font.attrib:
                        theme["fonts"]["body"] = minor_font.attrib['typeface']
                
                # Extract accent colors
                color_scheme = theme_element.find('.//a:clrScheme',
                    {'a': 'http://schemas.openxmlformats.org/drawingml/2006/main'})
                if color_scheme is not None:
                    for i in range(1, 7):
                        accent = color_scheme.find(f'.//a:accent{i}/a:srgbClr',
                            {'a': 'http://schemas.openxmlformats.org/drawingml/2006/main'})
                        if accent is not None and 'val' in accent.attrib:
                            theme["colors"][f"accent{i}"] = f"#{accent.attrib['val']}"
        except Exception:
            # Use defaults if theme parsing fails
            pass
        
        # Ensure we have all accent colors
        for i in range(1, 7):
            if f"accent{i}" not in theme["colors"]:
                # Default accent colors
                defaults = ["#4472C4", "#ED7D31", "#A5A5A5", "#FFC000", "#5B9BD5", "#70AD47"]
                theme["colors"][f"accent{i}"] = defaults[i-1]
        
        return theme
    
    def _extract_layouts(self) -> List[Dict[str, Any]]:
        """Extract layout information from slide layouts"""
        layouts = []
        
        for idx, layout in enumerate(self.prs.slide_layouts):
            layout_info = {
                "layout_id": f"layout_{idx}",
                "name": layout.name or f"Layout {idx + 1}",
                "placeholders": {
                    "title": False,
                    "bodies": 0,
                    "pictures": 0,
                    "table": False
                },
                "bbox": {}
            }
            
            for shape in layout.placeholders:
                if shape.placeholder_format.type == PP_PLACEHOLDER.TITLE:
                    layout_info["placeholders"]["title"] = True
                elif shape.placeholder_format.type in [PP_PLACEHOLDER.BODY, PP_PLACEHOLDER.OBJECT]:
                    layout_info["placeholders"]["bodies"] += 1
                elif shape.placeholder_format.type == PP_PLACEHOLDER.PICTURE:
                    layout_info["placeholders"]["pictures"] += 1
                elif shape.placeholder_format.type == PP_PLACEHOLDER.TABLE:
                    layout_info["placeholders"]["table"] = True
            
            layouts.append(layout_info)
        
        return layouts