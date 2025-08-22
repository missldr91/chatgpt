"""PPTX source document parser"""
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from typing import Dict, List, Any
import uuid
from pathlib import Path

class PPTXParser:
    """Parse source PPTX to extract page signatures"""
    
    def __init__(self, file_path: Path, template_id: str):
        self.prs = Presentation(str(file_path))
        self.template_id = template_id
        self.source_id = str(uuid.uuid4())
    
    def parse(self) -> Dict[str, Any]:
        """Extract signatures from all slides"""
        pages = []
        
        for idx, slide in enumerate(self.prs.slides):
            signature = self._analyze_slide(slide)
            pages.append({
                "idx": idx,
                "signature": signature,
                "warnings": []
            })
        
        return {
            "source_id": self.source_id,
            "type": "pptx",
            "pages": pages
        }
    
    def _analyze_slide(self, slide) -> Dict[str, Any]:
        """Analyze a single slide to extract its signature"""
        signature = {
            "title": False,
            "bullets": 0,
            "columns": 1,
            "images": 0,
            "table": False,
            "coverage": {"image": 0.0, "text": 0.0}
        }
        
        text_shapes = []
        image_area = 0
        text_area = 0
        total_area = float(self.prs.slide_width * self.prs.slide_height)
        
        # Find title (largest font in top 20% of slide)
        max_font_size = 0
        title_y_threshold = self.prs.slide_height * 0.2
        
        for shape in slide.shapes:
            try:
                if shape.has_text_frame:
                    # Calculate shape area
                    shape_area = shape.width * shape.height if hasattr(shape, 'width') else 0
                    text_area += shape_area
                    
                    # Check for title
                    if hasattr(shape, 'top') and shape.top < title_y_threshold:
                        for paragraph in shape.text_frame.paragraphs:
                            if paragraph.runs:
                                font_size = paragraph.runs[0].font.size
                                if font_size and font_size > max_font_size:
                                    max_font_size = font_size
                                    signature["title"] = True
                    
                    # Count bullets
                    for paragraph in shape.text_frame.paragraphs:
                        if paragraph.level > 0 or self._has_bullet(paragraph):
                            signature["bullets"] += 1
                    
                    # Store text shapes for column detection
                    if hasattr(shape, 'left') and hasattr(shape, 'width'):
                        text_shapes.append({
                            'left': shape.left,
                            'width': shape.width,
                            'centroid': shape.left + shape.width / 2
                        })
                
                # Count images
                elif shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    signature["images"] += 1
                    if hasattr(shape, 'width') and hasattr(shape, 'height'):
                        image_area += shape.width * shape.height
                
                # Check for tables
                elif shape.has_table:
                    signature["table"] = True
                    
            except Exception:
                continue
        
        # Detect columns (if text shapes are separated by > 25% of slide width)
        if len(text_shapes) >= 2:
            text_shapes.sort(key=lambda x: x['centroid'])
            for i in range(len(text_shapes) - 1):
                gap = text_shapes[i+1]['left'] - (text_shapes[i]['left'] + text_shapes[i]['width'])
                if gap > self.prs.slide_width * 0.25:
                    signature["columns"] = 2
                    break
        
        # Calculate coverage
        if total_area > 0:
            signature["coverage"]["image"] = min(1.0, image_area / total_area)
            signature["coverage"]["text"] = min(1.0, text_area / total_area)
        
        return signature
    
    def _has_bullet(self, paragraph) -> bool:
        """Check if paragraph has a bullet"""
        try:
            # Check for bullet character
            if paragraph.text and paragraph.text.strip():
                first_char = paragraph.text.strip()[0]
                if first_char in ['•', '·', '◦', '▪', '▫', '-', '–', '—', '*']:
                    return True
            
            # Check paragraph format
            pf = paragraph._element.find('.//a:buChar', 
                {'a': 'http://schemas.openxmlformats.org/drawingml/2006/main'})
            if pf is not None:
                return True
                
        except Exception:
            pass
        
        return False