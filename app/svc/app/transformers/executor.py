"""Transformation executor for generating final PPTX"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from typing import Dict, List, Any, Optional
from pathlib import Path
import uuid
import json
from PIL import Image, ImageDraw, ImageFont
import io
import base64

class TransformationExecutor:
    """Execute transformation plan to generate output PPTX"""
    
    def __init__(self, plan: Dict, template_path: Path, source_path: Path):
        self.plan = plan
        self.template_prs = Presentation(str(template_path))
        self.source_prs = Presentation(str(source_path))
        self.job_id = str(uuid.uuid4())
        self.output_prs = Presentation(str(template_path))
        
        # Clear template slides
        while len(self.output_prs.slides) > 0:
            rId = self.output_prs.slides._sldIdLst[0].rId
            self.output_prs.part.drop_rel(rId)
            del self.output_prs.slides._sldIdLst[0]
    
    def execute(self) -> Dict[str, Any]:
        """Execute the transformation plan"""
        preview_pngs = []
        report = {
            "greens": 0,
            "yellows": 0,
            "issues_by_type": {}
        }
        
        for slide_plan in self.plan["slides"]:
            idx = slide_plan["idx"]
            layout_id = slide_plan["chosen_layout_id"]
            issues = slide_plan["issues"]
            
            # Get source slide
            if idx < len(self.source_prs.slides):
                source_slide = self.source_prs.slides[idx]
                
                # Find template layout
                layout_idx = int(layout_id.split("_")[1])
                if layout_idx < len(self.template_prs.slide_layouts):
                    layout = self.template_prs.slide_layouts[layout_idx]
                else:
                    layout = self.template_prs.slide_layouts[0]
                
                # Create new slide with template layout
                new_slide = self.output_prs.slides.add_slide(layout)
                
                # Recompose content
                overflow_text = self._recompose_slide(source_slide, new_slide, layout)
                
                # Handle overflow by creating continuation slides
                if overflow_text:
                    issues.append("overflow")
                    self._create_continuation_slides(overflow_text, layout)
                
                # Generate preview PNG (simplified)
                preview = self._generate_preview(new_slide, issues)
                preview_pngs.append(preview)
            
            # Update report
            if not issues:
                report["greens"] += 1
            else:
                report["yellows"] += 1
                for issue in issues:
                    if issue not in report["issues_by_type"]:
                        report["issues_by_type"][issue] = 0
                    report["issues_by_type"][issue] += 1
        
        # Save output
        output_path = Path(f"data/jobs/{self.job_id}/output.pptx")
        output_path.parent.mkdir(parents=True, exist_ok=True)
        self.output_prs.save(str(output_path))
        
        return {
            "job_id": self.job_id,
            "status": "done",
            "artifact_url": f"/jobs/{self.job_id}/output.pptx",
            "preview_pngs": preview_pngs,
            "report": report
        }
    
    def _recompose_slide(self, source_slide, new_slide, layout) -> Optional[str]:
        """Recompose content from source to new slide"""
        overflow_text = ""
        
        # Extract text content from source
        title_text = ""
        body_texts = []
        
        for shape in source_slide.shapes:
            if shape.has_text_frame:
                text = shape.text_frame.text.strip()
                if text:
                    # Assume first text is title if in top portion
                    if not title_text and hasattr(shape, 'top'):
                        if shape.top < source_slide.part.slide_height * 0.3:
                            title_text = text
                        else:
                            body_texts.append(text)
                    else:
                        body_texts.append(text)
        
        # Fill template placeholders
        for shape in new_slide.placeholders:
            try:
                # Fill title
                if shape.placeholder_format.idx == 0:  # Title placeholder
                    if title_text:
                        shape.text = title_text
                
                # Fill body
                elif shape.placeholder_format.idx == 1:  # Content placeholder
                    if body_texts:
                        combined_text = "\n".join(body_texts)
                        
                        # Simple overflow detection
                        max_chars = 500  # Rough estimate
                        if len(combined_text) > max_chars:
                            shape.text = combined_text[:max_chars]
                            overflow_text = combined_text[max_chars:]
                        else:
                            shape.text = combined_text
            except Exception:
                continue
        
        return overflow_text if overflow_text else None
    
    def _create_continuation_slides(self, overflow_text: str, layout):
        """Create continuation slides for overflow content"""
        cont_slide = self.output_prs.slides.add_slide(layout)
        
        # Add (cont.) to title
        for shape in cont_slide.placeholders:
            try:
                if shape.placeholder_format.idx == 0:  # Title
                    shape.text = "(continued)"
                elif shape.placeholder_format.idx == 1:  # Content
                    shape.text = overflow_text
                    break
            except Exception:
                continue
    
    def _generate_preview(self, slide, issues: List[str]) -> str:
        """Generate a simple preview PNG as base64"""
        # Create a simple preview image
        img = Image.new('RGB', (400, 300), color='white')
        draw = ImageDraw.Draw(img)
        
        # Draw slide number and issues
        try:
            font = ImageFont.load_default()
        except:
            font = None
        
        # Draw title placeholder
        draw.rectangle([10, 10, 390, 50], outline='gray', width=2)
        draw.text((20, 20), "Title", fill='black', font=font)
        
        # Draw content placeholder
        draw.rectangle([10, 60, 390, 250], outline='gray', width=1)
        draw.text((20, 70), "Content", fill='gray', font=font)
        
        # Draw issues badge
        if issues:
            badge_color = 'yellow' if 'low_fit' not in issues else 'orange'
            draw.rectangle([350, 260, 390, 290], fill=badge_color)
            draw.text((355, 265), str(len(issues)), fill='black', font=font)
        else:
            draw.rectangle([350, 260, 390, 290], fill='green')
            draw.text((360, 265), "OK", fill='white', font=font)
        
        # Convert to base64
        buffer = io.BytesIO()
        img.save(buffer, format='PNG')
        img_str = base64.b64encode(buffer.getvalue()).decode()
        
        return f"data:image/png;base64,{img_str}"