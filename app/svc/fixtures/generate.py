"""Generate sample fixtures for tests and manual usage."""
from pptx import Presentation
from pptx.util import Inches
from pathlib import Path
from reportlab.pdfgen import canvas

DATA_DIR = Path(__file__).resolve().parent.parent / 'data' / 'fixtures'
DATA_DIR.mkdir(parents=True, exist_ok=True)

def build_template(path: Path):
    prs = Presentation()
    prs.slide_layouts[0].name = 'Title'
    prs.slide_layouts[1].name = 'Title and Content'
    prs.slide_layouts[2].name = 'Section Header'
    prs.slide_layouts[3].name = 'Two Content'
    prs.slide_layouts[4].name = 'Content with Caption'
    prs.save(path)


def build_source_pptx(path: Path):
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[2])
    slide.shapes.title.text = 'Section'
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = 'Bullets'
    tf = slide.shapes.placeholders[1].text_frame
    for i in range(3):
        p = tf.add_paragraph() if i>0 else tf.paragraphs[0]
        p.text = f'Point {i+1}'
        p.level = 1
    slide = prs.slides.add_slide(prs.slide_layouts[3])
    slide.shapes.title.text = 'Two Column'
    slide.shapes.placeholders[1].text = 'Left'
    slide.shapes.placeholders[2].text = 'Right'
    slide = prs.slides.add_slide(prs.slide_layouts[4])
    slide.shapes.title.text = 'Image'
    pic_path = DATA_DIR / 'python.png'
    if not pic_path.exists():
        from PIL import Image, ImageDraw
        img = Image.new('RGB', (200,200), 'blue')
        ImageDraw.Draw(img).text((50,90), 'IMG', fill='white')
        img.save(pic_path)
    slide.shapes.add_picture(str(pic_path), Inches(1), Inches(1))
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = 'Table'
    table = slide.shapes.add_table(2,2, Inches(1), Inches(2), Inches(6), Inches(1.5)).table
    table.cell(0,0).text = 'A'
    table.cell(0,1).text = 'B'
    table.cell(1,0).text = 'C'
    table.cell(1,1).text = 'D'
    prs.save(path)


def build_pdf(path: Path):
    c = canvas.Canvas(str(path))
    for i in range(5):
        c.drawString(72, 720, f"Page {i+1} - Sample text")
        c.showPage()
    c.save()

if __name__ == '__main__':
    build_template(DATA_DIR / 'brand_simple.pptx')
    build_source_pptx(DATA_DIR / 'mini_5slide.pptx')
    build_pdf(DATA_DIR / 'mini_pdf_5page.pdf')
