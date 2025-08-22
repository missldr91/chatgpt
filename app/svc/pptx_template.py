from pptx import Presentation
from pptx.enum.shapes import PP_PLACEHOLDER
from pathlib import Path
from typing import Dict, Any, List


def parse_template(path: Path) -> Dict[str, Any]:
    prs = Presentation(path)
    theme = getattr(prs.slide_master.part, 'theme', None)
    theme_meta = {
        'fonts': {
            'title': getattr(getattr(theme, 'major_font', None), 'latin', None).typeface if theme else 'Calibri',
            'body': getattr(getattr(theme, 'minor_font', None), 'latin', None).typeface if theme else 'Calibri',
        },
        'colors': {},
        'page_size': {'w': prs.slide_width, 'h': prs.slide_height}
    }
    if theme:
        color_scheme = getattr(theme, 'color_scheme', None)
        for i in range(1,7):
            if color_scheme:
                color = getattr(color_scheme, f'accent{i}')
                theme_meta['colors'][f'accent{i}'] = str(color.rgb)
    layout_catalog: List[Dict[str, Any]] = []
    for idx, layout in enumerate(prs.slide_layouts):
        placeholders = {'title': False, 'bodies':0, 'pictures':0, 'table':False}
        for shp in layout.placeholders:
            ph_type = shp.placeholder_format.type
            if ph_type == PP_PLACEHOLDER.TITLE:
                placeholders['title'] = True
            elif ph_type == PP_PLACEHOLDER.BODY:
                placeholders['bodies'] += 1
            elif ph_type == PP_PLACEHOLDER.PICTURE:
                placeholders['pictures'] += 1
            if shp.has_table:
                placeholders['table'] = True
        layout_catalog.append({
            'layout_id': idx,
            'name': layout.name,
            'placeholders': placeholders,
            'bbox': {}
        })
    return {'theme_meta': theme_meta, 'layout_catalog': layout_catalog}
