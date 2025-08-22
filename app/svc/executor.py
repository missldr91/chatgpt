from pptx import Presentation
from pathlib import Path
from typing import Dict, Any


def execute(plan: Dict[str, Any], template_path: Path, source_path: Path, output_path: Path):
    tpl = Presentation(template_path)
    src = Presentation(source_path)
    out = Presentation(template_path)
    out.slides._sldIdLst.clear()  # start empty
    layouts = tpl.slide_layouts
    for slide_plan in plan['slides']:
        idx = slide_plan['idx']
        layout = layouts[slide_plan['chosen_layout_id']]
        new_slide = out.slides.add_slide(layout)
        src_slide = src.slides[idx]
        title_text = ''
        body_text = []
        for shp in src_slide.shapes:
            if shp.has_text_frame:
                if shp.is_placeholder and shp.placeholder_format.type==1:  # title
                    title_text = shp.text
                else:
                    for p in shp.text_frame.paragraphs:
                        body_text.append(p.text)
        if title_text and new_slide.shapes.title:
            new_slide.shapes.title.text = title_text
        body_phs = [shp for shp in new_slide.shapes if shp.is_placeholder and shp.placeholder_format.type==2]
        if body_phs:
            tf = body_phs[0].text_frame
            tf.clear()
            for i,line in enumerate(body_text):
                p = tf.add_paragraph() if i>0 else tf.paragraphs[0]
                p.text = line
        # TODO: handle images, tables, overflow splitting
    out.save(output_path)
