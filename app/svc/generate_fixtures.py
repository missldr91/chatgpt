"""Generate sample fixture files for testing"""
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from reportlab.lib.pagesizes import letter
from reportlab.pdfgen import canvas
from reportlab.lib.utils import ImageReader
from PIL import Image, ImageDraw
import io

def create_brand_template():
    """Create a simple brand template with 5 layouts"""
    prs = Presentation()
    
    # Customize slide size (16:9)
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(5.625)
    
    print("Creating brand template with 5 layouts...")
    
    # The default template already has layouts, we'll use them
    # and just save with our customizations
    
    output_path = Path("data/fixtures/templates/brand_simple.pptx")
    output_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(output_path))
    
    print(f"✓ Created: {output_path}")
    return output_path

def create_sample_source():
    """Create a 5-slide sample source presentation"""
    prs = Presentation()
    
    print("Creating 5-slide sample source...")
    
    # Slide 1: Section Header
    slide1 = prs.slides.add_slide(prs.slide_layouts[2])  # Section Header
    slide1.shapes.title.text = "Introduction"
    if len(slide1.placeholders) > 1:
        slide1.placeholders[1].text = "Welcome to our presentation"
    
    # Slide 2: Title + Bullets
    slide2 = prs.slides.add_slide(prs.slide_layouts[1])  # Title and Content
    slide2.shapes.title.text = "Key Points"
    content = slide2.placeholders[1].text_frame
    content.text = "First main point about our product"
    p = content.add_paragraph()
    p.text = "Second important feature"
    p.level = 1
    p = content.add_paragraph()
    p.text = "Third key benefit"
    p.level = 1
    p = content.add_paragraph()
    p.text = "Additional value proposition"
    p.level = 2
    
    # Slide 3: Two Column
    slide3 = prs.slides.add_slide(prs.slide_layouts[3])  # Two Content
    slide3.shapes.title.text = "Comparison"
    # Add text to left placeholder
    for idx, placeholder in enumerate(slide3.placeholders):
        if idx == 1:  # Left content
            tf = placeholder.text_frame
            tf.text = "Option A Benefits:"
            p = tf.add_paragraph()
            p.text = "• Cost effective"
            p = tf.add_paragraph()
            p.text = "• Easy to implement"
            p = tf.add_paragraph()
            p.text = "• Proven results"
        elif idx == 2:  # Right content
            tf = placeholder.text_frame
            tf.text = "Option B Benefits:"
            p = tf.add_paragraph()
            p.text = "• More features"
            p = tf.add_paragraph()
            p.text = "• Better performance"
            p = tf.add_paragraph()
            p.text = "• Future proof"
    
    # Slide 4: Image Heavy (using shapes as placeholders)
    slide4 = prs.slides.add_slide(prs.slide_layouts[5])  # Blank layout
    title = slide4.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(1))
    title.text_frame.text = "Visual Overview"
    title.text_frame.paragraphs[0].font.size = Pt(32)
    title.text_frame.paragraphs[0].font.bold = True
    
    # Add placeholder rectangles to represent images
    left = Inches(0.5)
    top = Inches(1.5)
    for i in range(3):
        rect = slide4.shapes.add_shape(
            1,  # Rectangle
            left + (i * Inches(3.2)),
            top,
            Inches(3),
            Inches(2)
        )
        rect.fill.solid()
        rect.fill.fore_color.rgb = RGBColor(200, 200, 200)
        
        # Add caption
        caption = slide4.shapes.add_textbox(
            left + (i * Inches(3.2)),
            top + Inches(2.1),
            Inches(3),
            Inches(0.5)
        )
        caption.text_frame.text = f"Image {i+1}"
        caption.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # Slide 5: Table
    slide5 = prs.slides.add_slide(prs.slide_layouts[5])  # Blank layout
    title = slide5.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(1))
    title.text_frame.text = "Quarterly Results"
    title.text_frame.paragraphs[0].font.size = Pt(32)
    title.text_frame.paragraphs[0].font.bold = True
    
    # Add a table
    rows, cols = 5, 4
    left = Inches(1)
    top = Inches(1.5)
    width = Inches(8)
    height = Inches(3)
    
    table = slide5.shapes.add_table(rows, cols, left, top, width, height).table
    
    # Header row
    headers = ["Quarter", "Revenue", "Profit", "Growth"]
    for i, header in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = header
        cell.text_frame.paragraphs[0].font.bold = True
    
    # Data rows
    data = [
        ["Q1", "$1.2M", "$200K", "15%"],
        ["Q2", "$1.5M", "$300K", "25%"],
        ["Q3", "$1.8M", "$400K", "20%"],
        ["Q4", "$2.1M", "$500K", "17%"]
    ]
    
    for row_idx, row_data in enumerate(data, 1):
        for col_idx, value in enumerate(row_data):
            table.cell(row_idx, col_idx).text = value
    
    output_path = Path("data/fixtures/sources/mini_5slide.pptx")
    output_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(output_path))
    
    print(f"✓ Created: {output_path}")
    return output_path

def create_sample_pdf():
    """Create a 5-page sample PDF"""
    output_path = Path("data/fixtures/sources/mini_pdf_5page.pdf")
    output_path.parent.mkdir(parents=True, exist_ok=True)
    
    print("Creating 5-page sample PDF...")
    
    c = canvas.Canvas(str(output_path), pagesize=letter)
    width, height = letter
    
    # Page 1: Title page
    c.setFont("Helvetica-Bold", 36)
    c.drawString(100, height - 100, "Annual Report 2024")
    c.setFont("Helvetica", 18)
    c.drawString(100, height - 150, "Company Performance Overview")
    c.showPage()
    
    # Page 2: Bullet points
    c.setFont("Helvetica-Bold", 24)
    c.drawString(50, height - 50, "Executive Summary")
    c.setFont("Helvetica", 12)
    bullets = [
        "• Revenue increased by 45% year-over-year",
        "• Expanded to 3 new international markets",
        "• Launched 5 innovative products",
        "• Customer satisfaction rate: 94%",
        "• Team grew from 50 to 120 employees"
    ]
    y_pos = height - 100
    for bullet in bullets:
        c.drawString(70, y_pos, bullet)
        y_pos -= 25
    c.showPage()
    
    # Page 3: Two columns
    c.setFont("Helvetica-Bold", 24)
    c.drawString(50, height - 50, "Market Analysis")
    
    # Left column
    c.setFont("Helvetica-Bold", 14)
    c.drawString(50, height - 100, "Strengths")
    c.setFont("Helvetica", 11)
    left_items = [
        "- Market leader position",
        "- Strong brand recognition",
        "- Innovative technology",
        "- Excellent team"
    ]
    y_pos = height - 130
    for item in left_items:
        c.drawString(50, y_pos, item)
        y_pos -= 20
    
    # Right column
    c.setFont("Helvetica-Bold", 14)
    c.drawString(width/2 + 50, height - 100, "Opportunities")
    c.setFont("Helvetica", 11)
    right_items = [
        "- Emerging markets",
        "- Digital transformation",
        "- Strategic partnerships",
        "- New product lines"
    ]
    y_pos = height - 130
    for item in right_items:
        c.drawString(width/2 + 50, y_pos, item)
        y_pos -= 20
    c.showPage()
    
    # Page 4: Image placeholder
    c.setFont("Helvetica-Bold", 24)
    c.drawString(50, height - 50, "Product Showcase")
    
    # Draw rectangles to represent images
    c.setStrokeColorRGB(0.5, 0.5, 0.5)
    c.setFillColorRGB(0.9, 0.9, 0.9)
    c.rect(50, height - 300, 200, 150, fill=1)
    c.rect(300, height - 300, 200, 150, fill=1)
    c.rect(50, height - 500, 200, 150, fill=1)
    
    c.setFont("Helvetica", 10)
    c.drawString(125, height - 225, "[Product Image 1]")
    c.drawString(375, height - 225, "[Product Image 2]")
    c.drawString(125, height - 425, "[Product Image 3]")
    c.showPage()
    
    # Page 5: Table
    c.setFont("Helvetica-Bold", 24)
    c.drawString(50, height - 50, "Financial Summary")
    
    # Draw simple table
    c.setFont("Helvetica", 12)
    table_data = [
        ["Metric", "Q1", "Q2", "Q3", "Q4"],
        ["Revenue", "$2.1M", "$2.5M", "$2.8M", "$3.2M"],
        ["Profit", "$0.4M", "$0.5M", "$0.6M", "$0.8M"],
        ["Growth", "12%", "19%", "12%", "14%"],
        ["Customers", "1,200", "1,450", "1,680", "1,950"]
    ]
    
    y_start = height - 120
    x_positions = [50, 150, 220, 290, 360]
    
    for row_idx, row in enumerate(table_data):
        y_pos = y_start - (row_idx * 25)
        font = "Helvetica-Bold" if row_idx == 0 else "Helvetica"
        c.setFont(font, 12)
        for col_idx, cell in enumerate(row):
            c.drawString(x_positions[col_idx], y_pos, cell)
    
    c.save()
    print(f"✓ Created: {output_path}")
    return output_path

def main():
    """Generate all fixtures"""
    print("\n=== Generating Fixture Files ===\n")
    
    # Ensure directories exist
    Path("data/fixtures/templates").mkdir(parents=True, exist_ok=True)
    Path("data/fixtures/sources").mkdir(parents=True, exist_ok=True)
    
    # Generate files
    create_brand_template()
    create_sample_source()
    create_sample_pdf()
    
    print("\n✓ All fixtures generated successfully!")
    print("\nFixture files location:")
    print("  - data/fixtures/templates/brand_simple.pptx")
    print("  - data/fixtures/sources/mini_5slide.pptx")
    print("  - data/fixtures/sources/mini_pdf_5page.pdf")

if __name__ == "__main__":
    main()