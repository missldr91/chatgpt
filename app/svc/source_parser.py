from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pathlib import Path
from typing import Dict, Any, List
import pdfplumber


def parse_pptx(path: Path) -> List[Dict[str, Any]]:
    prs = Presentation(path)
    slides = []
    page_w, page_h = prs.slide_width, prs.slide_height
    for idx, slide in enumerate(prs.slides):
        title = False
        bullets = 0
        columns = 1
        images = 0
        table = False
        text_boxes = []
        for shp in slide.shapes:
            if shp.has_text_frame:
                text_boxes.append(shp)
                for p in shp.text_frame.paragraphs:
                    if p.level > 0 or p.text.strip().startswith(('•','-','–')):
                        bullets += 1
                if shp.top < page_h * 0.2:
                    title = True
            elif shp.shape_type == MSO_SHAPE_TYPE.PICTURE:
                images += 1
            if shp.has_table:
                table = True
        if len(text_boxes) >= 2:
            centers = [tb.left + tb.width/2 for tb in text_boxes]
            centers.sort()
            if centers[-1] - centers[0] > page_w*0.25:
                columns = 2
        coverage = {'image':0,'text':0}
        slide_area = page_w * page_h
        for tb in text_boxes:
            coverage['text'] += (tb.width*tb.height)/slide_area
        for shp in slide.shapes:
            if shp.shape_type == MSO_SHAPE_TYPE.PICTURE:
                coverage['image'] += (shp.width*shp.height)/slide_area
        slides.append({'idx': idx, 'signature': {'title': title, 'bullets': bullets, 'columns': columns, 'images': images, 'table': table, 'coverage': coverage}, 'warnings': []})
    return slides


def parse_pdf(path: Path) -> List[Dict[str, Any]]:
    slides = []
    with pdfplumber.open(path) as pdf:
        for i, page in enumerate(pdf.pages):
            text = page.extract_text() or ''
            lines = text.splitlines()
            title = bool(lines)
            bullets = sum(1 for line in lines if line.strip().startswith(('•','-','–')))
            images = len(page.images)
            coverage = {'image':0,'text':0}
            slides.append({'idx': i, 'signature': {'title': title, 'bullets': bullets, 'columns': 1, 'images': images, 'table': False, 'coverage': coverage}, 'warnings': []})
    return slides
