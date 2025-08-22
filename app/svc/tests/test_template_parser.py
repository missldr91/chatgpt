"""Tests for template parser"""
import pytest
from pathlib import Path
import sys
sys.path.append(str(Path(__file__).parent.parent))

from app.parsers import TemplateParser
from pptx import Presentation

def create_test_template():
    """Create a test template PPTX"""
    prs = Presentation()
    
    # Add some layouts (built-in layouts)
    # The default template has 11 layouts
    
    # Save test template
    test_path = Path("test_template.pptx")
    prs.save(str(test_path))
    return test_path

def test_template_parser_init():
    """Test template parser initialization"""
    test_path = create_test_template()
    try:
        parser = TemplateParser(test_path)
        assert parser.template_id is not None
        assert parser.prs is not None
    finally:
        test_path.unlink()

def test_template_parse():
    """Test template parsing"""
    test_path = create_test_template()
    try:
        parser = TemplateParser(test_path)
        result = parser.parse()
        
        # Check structure
        assert "template_id" in result
        assert "theme_meta" in result
        assert "layout_catalog" in result
        
        # Check theme meta
        theme = result["theme_meta"]
        assert "fonts" in theme
        assert "colors" in theme
        assert "page_size" in theme
        assert theme["fonts"]["title"] is not None
        assert theme["fonts"]["body"] is not None
        
        # Check colors
        for i in range(1, 7):
            assert f"accent{i}" in theme["colors"]
        
        # Check layouts
        layouts = result["layout_catalog"]
        assert len(layouts) > 0
        
        for layout in layouts:
            assert "layout_id" in layout
            assert "name" in layout
            assert "placeholders" in layout
            assert "title" in layout["placeholders"]
            assert "bodies" in layout["placeholders"]
            assert "pictures" in layout["placeholders"]
            
    finally:
        test_path.unlink()

def test_theme_extraction():
    """Test theme extraction specifically"""
    test_path = create_test_template()
    try:
        parser = TemplateParser(test_path)
        theme = parser._extract_theme()
        
        # Check default values are set
        assert theme["fonts"]["title"] in ["Calibri", "Arial", "Office Theme"]
        assert theme["fonts"]["body"] in ["Calibri", "Arial", "Office Theme"]
        assert theme["page_size"]["w"] > 0
        assert theme["page_size"]["h"] > 0
        
        # Check all accent colors exist
        for i in range(1, 7):
            assert f"accent{i}" in theme["colors"]
            assert theme["colors"][f"accent{i}"].startswith("#")
            
    finally:
        test_path.unlink()

def test_layout_extraction():
    """Test layout catalog extraction"""
    test_path = create_test_template()
    try:
        parser = TemplateParser(test_path)
        layouts = parser._extract_layouts()
        
        assert isinstance(layouts, list)
        assert len(layouts) > 0
        
        # Check first layout (usually Title Slide)
        first_layout = layouts[0]
        assert first_layout["layout_id"] == "layout_0"
        assert first_layout["name"] is not None
        assert isinstance(first_layout["placeholders"]["bodies"], int)
        assert isinstance(first_layout["placeholders"]["pictures"], int)
        
    finally:
        test_path.unlink()

if __name__ == "__main__":
    pytest.main([__file__, "-v"])