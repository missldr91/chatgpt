"""Tests for signature extraction"""
import pytest
from pathlib import Path
import sys
sys.path.append(str(Path(__file__).parent.parent))

from app.parsers import PPTXParser
from pptx import Presentation
from pptx.util import Inches, Pt

def create_test_source():
    """Create a test source PPTX with various slide types"""
    prs = Presentation()
    
    # Slide 1: Title slide
    slide1 = prs.slides.add_slide(prs.slide_layouts[0])
    slide1.shapes.title.text = "Test Presentation"
    slide1.placeholders[1].text = "Subtitle text"
    
    # Slide 2: Title and bullets
    slide2 = prs.slides.add_slide(prs.slide_layouts[1])
    slide2.shapes.title.text = "Bullet Points"
    content = slide2.placeholders[1].text_frame
    content.text = "First bullet"
    for i in range(2, 5):
        p = content.add_paragraph()
        p.text = f"Bullet {i}"
        p.level = 1
    
    # Slide 3: Two column layout (simulate)
    slide3 = prs.slides.add_slide(prs.slide_layouts[3])  # Two Content layout
    slide3.shapes.title.text = "Two Columns"
    
    # Save test source
    test_path = Path("test_source.pptx")
    prs.save(str(test_path))
    return test_path

def test_pptx_parser_init():
    """Test PPTX parser initialization"""
    test_path = create_test_source()
    try:
        parser = PPTXParser(test_path, "template_123")
        assert parser.source_id is not None
        assert parser.template_id == "template_123"
        assert parser.prs is not None
    finally:
        test_path.unlink()

def test_pptx_parse():
    """Test PPTX parsing"""
    test_path = create_test_source()
    try:
        parser = PPTXParser(test_path, "template_123")
        result = parser.parse()
        
        # Check structure
        assert "source_id" in result
        assert "type" in result
        assert "pages" in result
        assert result["type"] == "pptx"
        
        # Check pages
        pages = result["pages"]
        assert len(pages) >= 3
        
        # Check first page (title slide)
        page1 = pages[0]
        assert page1["idx"] == 0
        assert "signature" in page1
        assert page1["signature"]["title"] == True
        
        # Check second page (bullets)
        page2 = pages[1]
        assert page2["signature"]["title"] == True
        assert page2["signature"]["bullets"] >= 3
        
    finally:
        test_path.unlink()

def test_signature_analysis():
    """Test slide signature analysis"""
    test_path = create_test_source()
    try:
        parser = PPTXParser(test_path, "template_123")
        
        # Analyze first slide
        slide = parser.prs.slides[0]
        signature = parser._analyze_slide(slide)
        
        assert isinstance(signature, dict)
        assert "title" in signature
        assert "bullets" in signature
        assert "columns" in signature
        assert "images" in signature
        assert "table" in signature
        assert "coverage" in signature
        
        # Check coverage
        assert isinstance(signature["coverage"]["text"], float)
        assert isinstance(signature["coverage"]["image"], float)
        assert 0 <= signature["coverage"]["text"] <= 1
        assert 0 <= signature["coverage"]["image"] <= 1
        
    finally:
        test_path.unlink()

def test_bullet_detection():
    """Test bullet point detection"""
    test_path = create_test_source()
    try:
        parser = PPTXParser(test_path, "template_123")
        
        # Check slide with bullets
        slide = parser.prs.slides[1]
        signature = parser._analyze_slide(slide)
        
        assert signature["bullets"] > 0
        
    finally:
        test_path.unlink()

if __name__ == "__main__":
    pytest.main([__file__, "-v"])